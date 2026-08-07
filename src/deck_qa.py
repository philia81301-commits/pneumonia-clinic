# -*- coding: utf-8 -*-
"""
簡報版面稽核

這台機器沒有 LibreOffice，無法把投影片轉成圖片做目視 QA。
改以幾何與文字量估算做程式化檢查，涵蓋目視 QA 最常抓到的三類缺陷：
  1. 元件超出投影片邊界／邊界留白不足
  2. 文字超出容器（以 CJK 寬度估算）
  3. 文字方塊互相重疊

執行： PYTHONUTF8=1 python src/deck_qa.py output/台灣成人肺炎鏈球菌疫苗簡報.pptx
"""
import sys
from pptx import Presentation
from pptx.util import Emu

EMU_PER_IN = 914400
MARGIN_MIN = 0.5          # 英吋，邊界最小留白
OVERFLOW_TOL = 1.06       # 估算誤差容忍（6%）


def inches(v):
    return (v or 0) / EMU_PER_IN


def text_width_units(s, size_pt):
    """回傳文字寬度（點）。CJK 全形約 1 em，半形約 0.55 em。"""
    w = 0.0
    for ch in s:
        if ch == '\n':
            continue
        w += 1.0 if ord(ch) > 0x2E80 else 0.55
    return w * size_pt


def frame_needed_height(tf, box_w_in):
    """估算文字框所需高度（英吋）。"""
    usable_pt = max(box_w_in * 72.0 - 14, 20)   # 扣掉左右內距
    total_in = 0.0
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        if not runs:
            total_in += 0.12
            continue
        size = max([(r.font.size.pt if r.font.size else 18) for r in runs])
        line_h_in = size * 1.38 / 72.0
        for seg in ''.join(r.text for r in runs).split('\n'):
            w = text_width_units(seg, size)
            lines = max(1, int(w / usable_pt) + (1 if w % usable_pt else 0))
            total_in += lines * line_h_in
    return total_in


def collect(shape, out, prefix=''):
    if shape.shape_type == 6:  # GROUP
        for sh in shape.shapes:
            collect(sh, out, prefix)
        return
    out.append(shape)


# ── 出處稽核 ──
# 「資料來源很重要」：凡陳述事實數字或政策規定的投影片，都必須帶出處。
# 下列關鍵詞代表該頁有事實主張，必須在頁面文字中找到來源標記。
FACT_MARKERS = ('%', '歲', '價', '劑', '年', '週', 'PMID', 'CI', '例')
SOURCE_MARKERS = ('疾管署', '通函', '接種須知', '衛福部', '健保署', '內科學誌',
                  'PMID', 'Lancet', 'NEJM', 'MMWR', 'ACIP', 'Yellow Book',
                  'J Antimicrob', 'BMC', 'Hum Vaccin', 'Vaccine', 'JMII',
                  'Platt', 'Chiang', 'Huang', 'Pelton', 'clinicaltrials',
                  '查詢日', '出處', '來源')
# 純結構頁（封面、章節分隔、總覽、方法論宣告）不要求出處
EXEMPT_KEYWORDS = ('章節', '重點', '決策總覽', '只問三個問題', '來源位階',
                   '主要出處', '教學模組', '用語紅線', '破題', '類比')


def audit_sources(prs):
    """回傳缺少出處的投影片清單。"""
    missing = []
    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                texts.append(sh.text_frame.text)
            if getattr(sh, 'has_table', False):
                for r in sh.table.rows:
                    for c in r.cells:
                        if c.text.strip():
                            texts.append(c.text)
        blob = ' '.join(texts)
        if not blob.strip():
            continue
        # 章節封面／過場頁：文字量極少，不承載事實主張
        if len(blob.replace(' ', '')) < 40:
            continue
        if any(k in blob for k in EXEMPT_KEYWORDS):
            continue
        if not any(k in blob for k in FACT_MARKERS):
            continue
        if not any(k in blob for k in SOURCE_MARKERS):
            missing.append((idx, blob.replace('\n', ' ')[:46]))
    return missing


def main(path):
    prs = Presentation(path)
    SW, SH = inches(prs.slide_width), inches(prs.slide_height)
    print('投影片尺寸：%.2f × %.2f 英吋　共 %d 張\n' % (SW, SH, len(prs.slides.__iter__.__self__._sldIdLst)))

    issues = []
    for idx, slide in enumerate(prs.slides, 1):
        shapes = []
        for sh in slide.shapes:
            collect(sh, shapes)

        text_boxes = []
        for sh in shapes:
            x, y = inches(sh.left), inches(sh.top)
            w, h = inches(sh.width), inches(sh.height)

            # 1. 邊界
            if x < -0.01 or y < -0.01 or x + w > SW + 0.01 or y + h > SH + 0.01:
                issues.append((idx, 'OUT_OF_BOUNDS',
                               '形狀超出投影片：x=%.2f y=%.2f w=%.2f h=%.2f' % (x, y, w, h)))

            # 表格儲存格的文字溢出：GraphicFrame 沒有 text_frame，
            # 之前完全沒被檢查到，S3 的 IPD 說明文字因此被切掉。
            if getattr(sh, 'has_table', False):
                tbl = sh.table
                col_w = [c.width / EMU_PER_IN for c in tbl.columns]
                row_h = [r.height / EMU_PER_IN for r in tbl.rows]
                for ri, row in enumerate(tbl.rows):
                    for ci, cell in enumerate(row.cells):
                        txt_c = cell.text.strip()
                        if not txt_c:
                            continue
                        # 被合併吃掉的儲存格不重複檢查
                        if getattr(cell, 'is_spanned', False):
                            continue
                        # 跨列合併：實際可用高度是所跨各列之和
                        span = getattr(cell, 'span_height', 1) or 1
                        rh = sum(row_h[ri:ri + span])
                        size = 18
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                if r.font.size:
                                    size = max(size, r.font.size.pt)
                        if ci >= len(col_w):
                            continue
                        span_w = getattr(cell, 'span_width', 1) or 1
                        cw = sum(col_w[ci:ci + span_w])
                        usable = max(cw * 72.0 - 12, 20)
                        need = 0.0
                        for seg in txt_c.split('\n'):
                            wpt = text_width_units(seg, size)
                            lines = max(1, int(wpt / usable) + (1 if wpt % usable else 0))
                            need += lines * size * 1.35 / 72.0
                        if need > rh * OVERFLOW_TOL:
                            issues.append((idx, 'TABLE_CELL_OVERFLOW',
                                           '第 %d 列第 %d 欄需 %.2f" > 列高 %.2f"｜%s'
                                           % (ri + 1, ci + 1, need, rh, txt_c[:30].replace('\n', ' '))))
                continue

            if not sh.has_text_frame:
                continue
            txt = sh.text_frame.text.strip()
            if not txt:
                continue

            # 2. 文字溢框
            need = frame_needed_height(sh.text_frame, w)
            if need > h * OVERFLOW_TOL:
                issues.append((idx, 'TEXT_OVERFLOW',
                               '需 %.2f" > 框高 %.2f"｜%s' % (need, h, txt[:34].replace('\n', ' '))))

            # 3. 邊界留白（僅檢查含文字者）
            if x < MARGIN_MIN - 0.01 or x + w > SW - MARGIN_MIN + 0.01:
                issues.append((idx, 'TIGHT_MARGIN',
                               '左右留白不足：x=%.2f 右緣=%.2f｜%s' % (x, x + w, txt[:24].replace('\n', ' '))))

            text_boxes.append((x, y, w, h, txt))

        # 4. 文字方塊重疊
        for i in range(len(text_boxes)):
            for j in range(i + 1, len(text_boxes)):
                a, b = text_boxes[i], text_boxes[j]
                ox = min(a[0] + a[2], b[0] + b[2]) - max(a[0], b[0])
                oy = min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1])
                if ox > 0.12 and oy > 0.12:
                    issues.append((idx, 'TEXT_OVERLAP',
                                   '「%s」×「%s」重疊 %.2f×%.2f"' %
                                   (a[4][:16].replace('\n', ' '), b[4][:16].replace('\n', ' '), ox, oy)))

    missing_src = audit_sources(prs)

    if not issues and not missing_src:
        print('✅ 版面稽核全部通過（邊界／溢框／留白／重疊）')
        print('✅ 出處稽核全部通過（每張有事實主張的投影片都帶來源）')
        return 0

    if missing_src:
        print('【MISSING_SOURCE】%d 張投影片有事實主張但未標出處' % len(missing_src))
        for idx, txt in missing_src:
            print('  第 %2d 張　%s' % (idx, txt))
        print()

    if not issues:
        print('✅ 版面稽核通過（邊界／溢框／留白／重疊）')
        return 1

    by_kind = {}
    for it in issues:
        by_kind.setdefault(it[1], []).append(it)
    for kind in ('OUT_OF_BOUNDS', 'TEXT_OVERFLOW', 'TABLE_CELL_OVERFLOW', 'TEXT_OVERLAP', 'TIGHT_MARGIN'):
        rows = by_kind.get(kind, [])
        if not rows:
            continue
        print('【%s】%d 項' % (kind, len(rows)))
        for idx, _, msg in rows:
            print('  第 %2d 張　%s' % (idx, msg))
        print()
    return 1


if __name__ == '__main__':
    sys.exit(main(sys.argv[1] if len(sys.argv) > 1 else 'output/台灣成人肺炎鏈球菌疫苗簡報.pptx'))
